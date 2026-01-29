import streamlit as st
from openai import OpenAI
import json
import os
import PyPDF2
from docx import Document
from pptx import Presentation
import openpyxl
import csv
import moviepy.editor as mp
from tempfile import NamedTemporaryFile
import io
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import landscape, A4
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont

# --- КОНФИГУРАЦИЯ ---
MODEL_GPT = "gpt-4o"
MODEL_WHISPER = "whisper-1"

class QuizQuestion:
    def __init__(self, scenario, options, correct_option_id, explanation=""):
        self.scenario = scenario
        self.options = options
        self.correct_option_id = correct_option_id
        self.explanation = explanation

class Quiz:
    def __init__(self, questions):
        self.questions = questions

def get_client(api_key):
    return OpenAI(api_key=api_key)

# --- 1. ОБРАБОТКА ФАЙЛОВ ---
def process_file_to_text(uploaded_file, api_key):
    client = get_client(api_key)
    file_ext = uploaded_file.name.split('.')[-1].lower()
    text_content = ""

    try:
        # PDF
        if file_ext == 'pdf':
            pdf_reader = PyPDF2.PdfReader(uploaded_file)
            for page in pdf_reader.pages:
                text_content += page.extract_text() + "\n"
        
        # DOCX
        elif file_ext in ['docx', 'doc']:
            doc = Document(uploaded_file)
            text_content = "\n".join([para.text for para in doc.paragraphs])
        
        # TEXT
        elif file_ext == 'txt':
            text_content = uploaded_file.getvalue().decode("utf-8")
        
        # PPTX (PowerPoint)
        elif file_ext == 'pptx':
            prs = Presentation(uploaded_file)
            for slide_num, slide in enumerate(prs.slides, 1):
                text_content += f"\n--- Слайд {slide_num} ---\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_content += shape.text + "\n"
                    # Извлекаем текст из таблиц
                    if shape.has_table:
                        for row in shape.table.rows:
                            row_text = [cell.text for cell in row.cells]
                            text_content += " | ".join(row_text) + "\n"
        
        # XLSX (Excel)
        elif file_ext in ['xlsx', 'xls']:
            wb = openpyxl.load_workbook(uploaded_file, data_only=True)
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                text_content += f"\n--- Лист: {sheet_name} ---\n"
                for row in sheet.iter_rows(values_only=True):
                    row_text = [str(cell) if cell is not None else "" for cell in row]
                    if any(row_text):
                        text_content += " | ".join(row_text) + "\n"
        
        # CSV
        elif file_ext == 'csv':
            content = uploaded_file.getvalue().decode("utf-8")
            reader = csv.reader(content.splitlines())
            for row in reader:
                text_content += " | ".join(row) + "\n"
        
        # TSV
        elif file_ext == 'tsv':
            content = uploaded_file.getvalue().decode("utf-8")
            reader = csv.reader(content.splitlines(), delimiter='\t')
            for row in reader:
                text_content += " | ".join(row) + "\n"
        
        # ВИДЕО И АУДИО
        elif file_ext in ['mp4', 'mov', 'avi', 'mkv', 'mp3', 'wav', 'm4a', 'mpeg4', 'webm', 'wmv', 'flv', 'ogg', 'aac', 'wma', '3gp', 'mpeg', 'mpg']:
            with st.status("🎬 Обработка видео/аудио...", expanded=True) as status:
                status.write("1. Извлекаем аудиодорожку...")
                text_content = transcribe_audio_video(uploaded_file, client, status)
                status.update(label="✅ Готово!", state="complete", expanded=False)
        
        else:
            st.warning(f"⚠️ Формат .{file_ext} пока не поддерживается")
            return ""

    except Exception as e:
        st.error(f"❌ Ошибка обработки файла: {e}")
        return ""

    if not text_content or not text_content.strip():
        st.warning("⚠️ Нет текста. Возможно, файл пустой или содержит только изображения.")
    
    return text_content

def transcribe_audio_video(uploaded_file, client, status_container):
    try:
        # Сохраняем во временный файл
        suffix = f".{uploaded_file.name.split('.')[-1]}"
        with NamedTemporaryFile(delete=False, suffix=suffix) as tmp_video:
            tmp_video.write(uploaded_file.getvalue())
            tmp_video_path = tmp_video.name

        audio_path = tmp_video_path + "_converted.mp3"
        
        # Конвертация через MoviePy (требует FFMPEG)
        status_container.write("2. Конвертация в формат MP3 (32kbps)...")
        
        if suffix.lower() in ['.mp4', '.mov', '.avi', '.mkv', '.webm', '.wmv', '.mpeg4']:
            video = mp.VideoFileClip(tmp_video_path)
            # Если видео длиннее 20 минут - режем
            if video.duration > 1200: 
                status_container.warning(f"Видео длинное ({int(video.duration)}с). Берем первые 20 мин.")
                video = video.subclip(0, 1200)
            
            video.audio.write_audiofile(audio_path, bitrate="32k", logger=None)
            video.close()
        else:
            # Аудио тоже прогоняем через конвертер для сжатия
            audio_clip = mp.AudioFileClip(tmp_video_path)
            audio_clip.write_audiofile(audio_path, bitrate="32k", logger=None)
            audio_clip.close()

        # Проверка размера
        size_mb = os.path.getsize(audio_path) / (1024*1024)
        status_container.write(f"3. Отправка в Whisper AI ({size_mb:.1f} MB)...")

        if size_mb > 24:
            st.error("Файл слишком большой (>25MB) даже после сжатия.")
            return ""

        with open(audio_path, "rb") as audio_file:
            transcript = client.audio.transcriptions.create(
                model=MODEL_WHISPER, file=audio_file, response_format="text"
            )
        
        # Чистка
        try: os.remove(tmp_video_path); os.remove(audio_path)
        except: pass
        
        return transcript

    except Exception as e:
        st.error(f"Ошибка транскрибации (FFMPEG/Whisper): {str(e)}")
        if "ffmpeg" in str(e).lower():
            st.error("🚨 На сервере не найден FFMPEG. Установите: sudo apt install ffmpeg")
        return ""

# --- 2. ГЕНЕРАЦИЯ ТЕСТА ---
def generate_quiz_ai(text, num_questions, difficulty, language):
    client = get_client(st.secrets["OPENAI_API_KEY"])
    if not text: return Quiz([])
    
    prompt = f"""
You are an expert quiz creator. Create an engaging quiz based on the following text.

TEXT:
{text[:25000]}

REQUIREMENTS:
- Language: {language}
- Difficulty: {difficulty}
- Number of questions: {num_questions}
- Each question must have exactly 4 options
- Include a brief explanation (1-2 sentences) for why the correct answer is right

OUTPUT FORMAT (strict JSON):
{{
  "questions": [
    {{
      "scenario": "Question text here?",
      "options": ["Option A", "Option B", "Option C", "Option D"],
      "correct_option_id": 0,
      "explanation": "Brief explanation why this answer is correct."
    }}
  ]
}}
"""
    try:
        response = client.chat.completions.create(
            model=MODEL_GPT,
            messages=[{"role": "user", "content": prompt}],
            response_format={"type": "json_object"}
        )
        data = json.loads(response.choices[0].message.content)
        return Quiz([QuizQuestion(q['scenario'], q['options'], q['correct_option_id'], q.get('explanation', '')) for q in data['questions']])
    except Exception as e: return Quiz([QuizQuestion(f"Error: {e}", ["OK"], 0)])

def generate_methodologist_hints(text, language):
    if not text: return "Нет текста."
    client = get_client(st.secrets["OPENAI_API_KEY"])
    try:
        res = client.chat.completions.create(
            model=MODEL_GPT, messages=[{"role": "user", "content": f"3 learning tips for: {text[:5000]}. Lang: {language}"}]
        )
        return res.choices[0].message.content
    except: return "Советы недоступны."

# --- 3. ЭКСПОРТ ---
def create_html_quiz(quiz_obj, filename):
    js_data = []
    for q in quiz_obj.questions:
        js_data.append({"question": q.scenario, "options": q.options, "correct": q.correct_option_id})
    json_str = json.dumps(js_data)
    
    html = f"""
    <!DOCTYPE html>
    <html><head><meta charset='UTF-8'><style>body{{font-family:sans-serif;padding:20px;max-width:800px;margin:0 auto}} .card{{border:1px solid #ccc;padding:15px;margin-bottom:15px;border-radius:8px}} .btn{{background:#007bff;color:white;padding:10px 20px;border:none;cursor:pointer}} .correct{{color:green;font-weight:bold}} .wrong{{color:red;font-weight:bold}}</style></head>
    <body><h1>Test: {filename}</h1><div id="q"></div><button class="btn" onclick="check()">Check</button><h2 id="sc"></h2>
    <script>const d={json_str}; function r(){{let h='';d.forEach((q,i)=>{{h+=`<div class='card'><h3>${{i+1}}. ${{q.question}}</h3>`;q.options.forEach((o,j)=>{{h+=`<label style='display:block'><input type='radio' name='q${{i}}' value='${{j}}'> ${{o}}</label>`}});h+=`<div id='r${{i}}'></div></div>`}});document.getElementById('q').innerHTML=h}} r();
    function check(){{let s=0;d.forEach((q,i)=>{{let el=document.querySelector(`input[name='q${{i}}']:checked`);let r=document.getElementById(`r${{i}}`);if(el&&parseInt(el.value)===q.correct){{s++;r.innerHTML="<span class='correct'>OK</span>"}}else{{r.innerHTML=`<span class='wrong'>Wrong. Answer: ${{q.options[q.correct]}}</span>`}}}});document.getElementById('sc').innerText=`Score: ${{s}}/${{d.length}}`}}</script></body></html>
    """
    return html.encode('utf-8')

def remove_white_background(img):
    if img.mode != "RGBA":
        img = img.convert("RGBA")
    datas = img.getdata()
    new_data = []
    for item in datas:
        if item[0] > 240 and item[1] > 240 and item[2] > 240:
            new_data.append((255, 255, 255, 0))
        else:
            new_data.append(item)
    img.putdata(new_data)
    return img

def create_certificate(student_name, course_name, logo_file=None, signature_file=None):
    from reportlab.lib.colors import HexColor
    from reportlab.lib.utils import ImageReader
    from PIL import Image
    from datetime import datetime
    import random
    buffer = io.BytesIO()
    c = canvas.Canvas(buffer, pagesize=landscape(A4))
    width, height = landscape(A4)
    bg = HexColor("#0E1117")
    neon = HexColor("#00D4FF")
    purple = HexColor("#7D3CFF")
    light = HexColor("#FAFAFA")
    muted = HexColor("#979797")
    c.setFillColor(bg)
    c.rect(0, 0, width, height, fill=True, stroke=False)
    c.setStrokeColor(neon)
    c.setLineWidth(3)
    c.rect(25, 25, width-50, height-50)
    c.setStrokeColor(purple)
    c.setLineWidth(1)
    c.rect(35, 35, width-70, height-70)
    if logo_file:
        try:
            logo_file.seek(0)
            img = Image.open(io.BytesIO(logo_file.getvalue()))
            img = remove_white_background(img)
            r = min(250/img.width, 120/img.height)
            nw, nh = int(img.width*r), int(img.height*r)
            buf = io.BytesIO()
            img.save(buf, format="PNG")
            buf.seek(0)
            c.drawImage(ImageReader(buf), 50, height-nh-50, width=nw, height=nh, mask="auto")
        except: pass
    c.setFillColor(neon)
    c.setFont("Helvetica-Bold", 12)
    c.drawRightString(width-50, height-60, "VYUD AI CERTIFIED")
    c.setFillColor(light)
    c.setFont("Helvetica-Bold", 48)
    c.drawCentredString(width/2, height-180, "CERTIFICATE")
    c.setFillColor(neon)
    c.setFont("Helvetica", 20)
    c.drawCentredString(width/2, height-220, "OF COMPLETION")
    c.setStrokeColor(purple)
    c.setLineWidth(2)
    c.line(width/2-150, height-245, width/2+150, height-245)
    c.setFillColor(muted)
    c.setFont("Helvetica", 14)
    c.drawCentredString(width/2, height-280, "This certifies that")
    c.setFillColor(light)
    c.setFont("Helvetica-Bold", 36)
    c.drawCentredString(width/2, height-330, student_name)
    c.setFillColor(muted)
    c.setFont("Helvetica", 14)
    c.drawCentredString(width/2, height-370, "has successfully completed")
    c.setFillColor(neon)
    c.setFont("Helvetica-BoldOblique", 24)
    c.drawCentredString(width/2, height-410, course_name)
    if signature_file:
        try:
            signature_file.seek(0)
            img = Image.open(io.BytesIO(signature_file.getvalue()))
            img = remove_white_background(img)
            r = min(150/img.width, 60/img.height)
            nw, nh = int(img.width*r), int(img.height*r)
            buf = io.BytesIO()
            img.save(buf, format="PNG")
            buf.seek(0)
            c.drawImage(ImageReader(buf), 100, 80, width=nw, height=nh, mask="auto")
            c.setStrokeColor(muted)
            c.line(80, 75, 250, 75)
            c.setFillColor(muted)
            c.setFont("Helvetica", 10)
            c.drawString(80, 60, "Authorized Signature")
        except: pass
    c.setFillColor(muted)
    c.setFont("Helvetica", 12)
    dt = datetime.now().strftime("%B %d, %Y")
    c.drawRightString(width-80, 80, f"Issued: {dt}")
    c.setFont("Helvetica", 10)
    cid = f"VYUD-{random.randint(10000,99999)}"
    c.drawRightString(width-80, 60, f"Certificate ID: {cid}")
    c.save()
    buffer.seek(0)
    return buffer.getvalue()

def transcribe_for_bot(file_path):
    """Транскрибация для Telegram бота - принимает путь к файлу"""
    try:
        import os
        from openai import OpenAI
        
        client = OpenAI(api_key=os.environ.get("OPENAI_API_KEY"))
        
        # Конвертация в mp3 (если это видео)
        audio_path = file_path + "_converted.mp3"
        
        if file_path.lower().endswith(('.mp4', '.mov', '.avi', '.mkv', '.webm')):
            video = mp.VideoFileClip(file_path)
            if video.duration > 1200:
                video = video.subclip(0, 1200)
            video.audio.write_audiofile(audio_path, bitrate="32k", logger=None)
            video.close()
        else:
            # Аудио - тоже конвертируем для сжатия
            audio_clip = mp.AudioFileClip(file_path)
            audio_clip.write_audiofile(audio_path, bitrate="32k", logger=None)
            audio_clip.close()
        
        # Проверка размера
        size_mb = os.path.getsize(audio_path) / (1024*1024)
        if size_mb > 24:
            return "Error: File too large"
        
        # Whisper
        with open(audio_path, "rb") as audio_file:
            transcript = client.audio.transcriptions.create(
                model="whisper-1", file=audio_file, response_format="text"
            )
        
        # Чистка
        try: 
            os.remove(file_path)
            os.remove(audio_path)
        except: pass
        
        return transcript
        
    except Exception as e:
        return f"Error: {str(e)}"


# === ВЕРСИЯ ДЛЯ БОТА (без Streamlit) ===
def process_file_to_text_bot(file_path, file_name, api_key):
    """Обработка документов для Telegram бота - принимает путь к файлу"""
    import os
    from openai import OpenAI
    
    client = OpenAI(api_key=api_key)
    file_ext = file_name.split('.')[-1].lower()
    text_content = ""

    try:
        # PDF
        if file_ext == 'pdf':
            with open(file_path, 'rb') as f:
                pdf_reader = PyPDF2.PdfReader(f)
                for page in pdf_reader.pages:
                    extracted = page.extract_text()
                    if extracted:
                        text_content += extracted + "\n"
        
        # DOCX
        elif file_ext in ['docx', 'doc']:
            doc = Document(file_path)
            text_content = "\n".join([para.text for para in doc.paragraphs])
        
        # TEXT
        elif file_ext == 'txt':
            with open(file_path, 'r', encoding='utf-8') as f:
                text_content = f.read()
        
        # PPTX (PowerPoint)
        elif file_ext == 'pptx':
            prs = Presentation(file_path)
            for slide_num, slide in enumerate(prs.slides, 1):
                text_content += f"\n--- Слайд {slide_num} ---\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_content += shape.text + "\n"
                    if shape.has_table:
                        for row in shape.table.rows:
                            row_text = [cell.text for cell in row.cells]
                            text_content += " | ".join(row_text) + "\n"
        
        # XLSX (Excel)
        elif file_ext in ['xlsx', 'xls']:
            wb = openpyxl.load_workbook(file_path, data_only=True)
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                text_content += f"\n--- Лист: {sheet_name} ---\n"
                for row in sheet.iter_rows(values_only=True):
                    row_text = [str(cell) if cell is not None else "" for cell in row]
                    if any(row_text):
                        text_content += " | ".join(row_text) + "\n"
        
        else:
            return f"Error: Формат .{file_ext} не поддерживается"

    except Exception as e:
        return f"Error: {str(e)}"

    if not text_content or not text_content.strip():
        return "Error: Нет текста в документе"
    
    return text_content
